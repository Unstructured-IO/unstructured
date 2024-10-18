"""Partitioner for PPTX documents.

PPTX files are PowerPoint 2007+ documents. These are XML-based and "open" (documented ISO standard),
unlike the `.ppt` format which was binary and proprietary.
"""

from __future__ import annotations

import io
from tempfile import SpooledTemporaryFile
from typing import IO, Any, Iterator, Protocol, Sequence

import pptx
from pptx.presentation import Presentation
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture
from pptx.shapes.shapetree import _BaseGroupShapes  # pyright: ignore [reportPrivateUsage]
from pptx.slide import Slide
from pptx.text.text import _Paragraph  # pyright: ignore [reportPrivateUsage]

from unstructured.chunking import add_chunking_strategy
from unstructured.common.html_table import HtmlTable, htmlify_matrix_of_cell_texts
from unstructured.documents.elements import (
    Element,
    ElementMetadata,
    EmailAddress,
    ListItem,
    NarrativeText,
    PageBreak,
    Table,
    Text,
    Title,
)
from unstructured.file_utils.model import FileType
from unstructured.partition.common.metadata import apply_metadata, get_last_modified_date
from unstructured.partition.text_type import (
    is_email_address,
    is_possible_narrative_text,
    is_possible_title,
)
from unstructured.partition.utils.constants import PartitionStrategy
from unstructured.utils import is_temp_file_path, lazyproperty

DETECTION_ORIGIN = "pptx"


def register_picture_partitioner(picture_partitioner: AbstractPicturePartitioner) -> None:
    """Specify a pluggable sub-partitioner to be used for partitioning PPTX images."""
    PptxPartitionerOptions.register_picture_partitioner(picture_partitioner)


# ================================================================================================
# PPTX DOMAIN MODEL DEFINITIONS
# ================================================================================================


class AbstractPicturePartitioner(Protocol):
    """Defines the interface for a pluggable sub-partitioner for PPTX Picture objects.

    A PPTX Picture object generally contains an image (e.g. JPG, PNG) but can also contain other
    media types like a video or sound file. The interface classmethod generates zero-or-more
    elements from the specified Picture object. If the media in the picture object is not supported
    then it will silently return without generating any elements.
    """

    @classmethod
    def iter_elements(cls, picture: Picture, opts: PptxPartitionerOptions) -> Iterator[Element]:
        """Generate document elements derived from `picture`, a PPTX Picture shape."""
        ...


# ================================================================================================
# PARTITIONER
# ================================================================================================


@apply_metadata(FileType.PPTX)
@add_chunking_strategy
def partition_pptx(
    filename: str | None = None,
    *,
    file: IO[bytes] | None = None,
    include_page_breaks: bool = True,
    include_slide_notes: bool | None = None,
    infer_table_structure: bool = True,
    starting_page_number: int = 1,
    strategy: str = PartitionStrategy.FAST,
    **kwargs: Any,
) -> list[Element]:
    """Partition PowerPoint document in .pptx format into its document elements.

    Parameters
    ----------
    filename
        A string defining the target filename path.
    file
        A file-like object using "rb" mode --> open(filename, "rb").
    include_page_breaks
        If True, includes a PageBreak element between slides
    include_slide_notes
        If True, includes the slide notes as element
    infer_table_structure
        If True, any Table elements that are extracted will also have a metadata field
        named "text_as_html" where the table's text content is rendered into an html string.
        I.e., rows and cells are preserved.
        Whether True or False, the "text" field is always present in any Table element
        and is the text content of the table (no structure).
    starting_page_number
        Indicates what page number should be assigned to the first slide in the presentation.
        This information will be reflected in elements' metadata and can be be especially
        useful when partitioning a document that is part of a larger document.
    """
    opts = PptxPartitionerOptions(
        file=file,
        file_path=filename,
        include_page_breaks=include_page_breaks,
        include_slide_notes=include_slide_notes,
        infer_table_structure=infer_table_structure,
        strategy=strategy,
        starting_page_number=starting_page_number,
    )

    return list(_PptxPartitioner.iter_presentation_elements(opts))


class _PptxPartitioner:
    """Provides `.partition()` for PowerPoint 2007+ (.pptx) files."""

    def __init__(self, opts: PptxPartitionerOptions):
        self._opts = opts

    @classmethod
    def iter_presentation_elements(cls, opts: PptxPartitionerOptions) -> Iterator[Element]:
        """Partition MS Word documents (.docx format) into its document elements."""
        return cls(opts)._iter_presentation_elements()

    def _iter_presentation_elements(self) -> Iterator[Element]:
        """Generate each document-element in presentation in document order."""
        # -- This implementation composes a collection of iterators into a "combined" iterator
        # -- return value using `yield from`. You can think of the return value as an Element
        # -- stream and each `yield from` as "add elements found by this function to the stream".
        # -- This is functionally analogous to declaring `elements: List[Element] = []` at the top
        # -- and using `elements.extend()` for the results of each of the function calls, but is
        # -- more perfomant, uses less memory (avoids producing and then garbage-collecting all
        # -- those small lists), is more flexible for later iterator operations like filter,
        # -- chain, map, etc. and is perhaps more elegant and simpler to read once you have the
        # -- concept of what it's doing. You can see the same pattern repeating in the "sub"
        # -- functions like `._iter_shape_elements()` where the "just return when done"
        # -- characteristic of a generator avoids repeated code to form interim results into lists.

        for slide in self._presentation.slides:
            yield from self._opts.increment_page_number()
            yield from self._iter_maybe_slide_notes(slide)

            title_shape, shapes = self._order_shapes(slide)

            for shape in shapes:
                if shape.has_table:
                    assert isinstance(shape, GraphicFrame)
                    yield from self._iter_table_element(shape)
                elif shape.has_text_frame:
                    assert isinstance(shape, Shape)
                    if shape == title_shape:
                        yield from self._iter_title_shape_element(shape)
                    else:
                        yield from self._iter_shape_elements(shape)
                elif isinstance(shape, Picture):
                    yield from self._iter_picture_elements(shape)

                # -- otherwise ditch it, this would include charts, connectors (lines),
                # -- and free-form shapes (squiggly lines). Lines don't have text.

    def _is_bulleted_paragraph(self, paragraph: _Paragraph) -> bool:
        """True when `paragraph` has a bullet-charcter prefix.

        Bullet characters in the openxml schema are represented by buChar.
        """
        # -- True when XPath returns a non-empty list (nodeset) --
        return bool(paragraph._p.xpath("./a:pPr/a:buChar"))

    def _iter_maybe_slide_notes(self, slide: Slide) -> Iterator[NarrativeText]:
        """Generate zero-or-one NarrativeText element for the slide-notes."""
        # -- only emit slide-notes elements when enabled --
        if not self._opts.include_slide_notes:
            return

        # -- not all slides have a notes slide --
        if not slide.has_notes_slide:
            return

        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame

        # -- not all notes slides have a text-frame (it's created on first use) --
        if not notes_text_frame:
            return
        notes_text = notes_text_frame.text.strip()

        # -- not all notes text-frams contain text (if it's all deleted the text-frame remains) --
        if not notes_text:
            return

        yield NarrativeText(
            text=notes_text, metadata=self._opts.text_metadata(), detection_origin=DETECTION_ORIGIN
        )

    def _iter_picture_elements(self, picture: Picture) -> Iterator[Element]:
        """Generate elements derived from the image in `picture`."""
        # -- delegate this job to the pluggable Picture partitioner --
        PicturePartitionerCls = self._opts.picture_partitioner
        yield from PicturePartitionerCls.iter_elements(picture, self._opts)

    def _iter_shape_elements(self, shape: Shape) -> Iterator[Element]:
        """Generate Text or subtype element for each paragraph in `shape`."""
        if self._shape_is_off_slide(shape):
            return

        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text
            if text.strip() == "":
                continue

            level = paragraph.level or 0
            metadata = self._opts.text_metadata(category_depth=level)

            if self._is_bulleted_paragraph(paragraph):
                yield ListItem(text=text, metadata=metadata, detection_origin=DETECTION_ORIGIN)
            elif is_email_address(text):
                yield EmailAddress(text=text, detection_origin=DETECTION_ORIGIN)
            elif is_possible_narrative_text(text):
                yield NarrativeText(
                    text=text,
                    metadata=metadata,
                    detection_origin=DETECTION_ORIGIN,
                )
            elif is_possible_title(text):
                # If text is a title but not the title shape increment the category depth)
                metadata = self._opts.text_metadata(category_depth=level + 1)
                yield Title(text=text, metadata=metadata, detection_origin=DETECTION_ORIGIN)
            else:
                yield Text(text=text, metadata=metadata, detection_origin=DETECTION_ORIGIN)

    def _iter_table_element(self, graphfrm: GraphicFrame) -> Iterator[Table]:
        """Generate zero-or-one Table element for the table in `shape`.

        An empty table does not produce an element.
        """
        if not (rows := list(graphfrm.table.rows)):
            return

        html_text = htmlify_matrix_of_cell_texts(
            [[cell.text for cell in row.cells] for row in rows]
        )
        html_table = HtmlTable.from_html_text(html_text)

        if not html_table.text:
            return

        metadata = self._opts.table_metadata(
            html_table.html if self._opts.infer_table_structure else None
        )

        yield Table(text=html_table.text, metadata=metadata, detection_origin=DETECTION_ORIGIN)

    def _iter_title_shape_element(self, shape: Shape) -> Iterator[Element]:
        """Generate Title element for each paragraph in title `shape`.

        Text is most likely a title, but in the rare case that the title shape was used
        for the slide body text, also check for bulleted paragraphs."""
        if self._shape_is_off_slide(shape):
            return

        depth = 0
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text
            if text.strip() == "":
                continue

            if self._is_bulleted_paragraph(paragraph):
                bullet_depth = paragraph.level or 0
                yield ListItem(
                    text=text,
                    metadata=self._opts.text_metadata(category_depth=bullet_depth),
                    detection_origin=DETECTION_ORIGIN,
                )
            elif is_email_address(text):
                yield EmailAddress(text=text, detection_origin=DETECTION_ORIGIN)
            else:
                # increment the category depth by the paragraph increment in the shape
                yield Title(
                    text=text,
                    metadata=self._opts.text_metadata(category_depth=depth),
                    detection_origin=DETECTION_ORIGIN,
                )
                depth += 1  # Cannot enumerate because we want to skip empty paragraphs

    def _order_shapes(self, slide: Slide) -> tuple[Shape | None, Sequence[BaseShape]]:
        """Orders the shapes on `slide` from top to bottom and left to right.

        Returns the title shape if it exists and the ordered shapes."""

        def iter_shapes(shapes: _BaseGroupShapes) -> Iterator[BaseShape]:
            for shape in shapes:
                if isinstance(shape, GroupShape):
                    yield from iter_shapes(shape.shapes)
                else:
                    yield shape

        def sort_key(shape: BaseShape) -> tuple[int, int]:
            return shape.top or 0, shape.left or 0

        return slide.shapes.title, sorted(iter_shapes(slide.shapes), key=sort_key)

    @lazyproperty
    def _presentation(self) -> Presentation:
        """The python-pptx `Presentation` object loaded from the provided source file."""
        return pptx.Presentation(self._opts.pptx_file)

    def _shape_is_off_slide(self, shape: Shape) -> bool:
        # NOTE(robinson) - avoid processing shapes that are not on the actual slide
        # NOTE - skip check if no top or left position (shape displayed top left)
        return bool((shape.top and shape.left) and (shape.top < 0 or shape.left < 0))


class PptxPartitionerOptions:
    """Encapsulates partitioning option validation, computation, and application of defaults."""

    _PicturePartitionerCls = None
    """Sub-partitioner used to partition PPTX Picture (Image) shapes.

    This value has module lifetime and is updated by calling the `register_picture_partitioner()`
    function defined in this module. The value sent to `register_picture_partitioner()` must be a
    pluggable sub-partitioner implementing the `AbstractPicturePartitioner` interface. After
    registration, all picture shapes in subsequent PPTX documents will be partitioned by the
    specified picture sub-partitioner.
    """

    def __init__(
        self,
        *,
        file: IO[bytes] | None,
        file_path: str | None,
        include_page_breaks: bool,
        include_slide_notes: bool | None,
        infer_table_structure: bool,
        strategy: str,
        starting_page_number: int = 1,
    ):
        self._file = file
        self._file_path = file_path
        self._include_page_breaks = include_page_breaks
        self._include_slide_notes = include_slide_notes
        self._infer_table_structure = infer_table_structure
        self._strategy = strategy
        # -- options object maintains page-number state --
        self._page_counter = starting_page_number - 1

    @classmethod
    def register_picture_partitioner(cls, picture_partitioner: AbstractPicturePartitioner):
        """Specify a pluggable sub-partitioner to be used for partitioning PPTX images."""
        cls._PicturePartitionerCls = picture_partitioner

    @lazyproperty
    def include_page_breaks(self) -> bool:
        """When True, include `PageBreak` elements in element-stream.

        Note that regardless of this setting, page-breaks are detected, and page-number is tracked
        and included in element metadata. Only the presence of distinct `PageBreak` elements (which
        contain no text) in the element stream is affected.
        """
        return self._include_page_breaks

    @lazyproperty
    def include_slide_notes(self) -> bool:
        """When True, also partition any text found in slide notes as part of each slide."""
        return False if self._include_slide_notes is None else self._include_slide_notes

    def increment_page_number(self) -> Iterator[PageBreak]:
        """Increment page-number by 1 and generate a PageBreak element if enabled."""
        self._page_counter += 1
        # -- no page-break before first page --
        if self._page_counter < 2:
            return
        # -- only emit page-breaks when enabled --
        if self._include_page_breaks:
            yield PageBreak(
                "",
                detection_origin=DETECTION_ORIGIN,
                metadata=ElementMetadata(
                    last_modified=self.last_modified, page_number=self.page_number - 1
                ),
            )

    @lazyproperty
    def infer_table_structure(self) -> bool:
        """True when partitioner should compute and apply `text_as_html` metadata for tables."""
        return self._infer_table_structure

    @lazyproperty
    def last_modified(self) -> str | None:
        """The best last-modified date available, None if no sources are available."""
        if not self._file_path:
            return None

        return (
            None if is_temp_file_path(self._file_path) else get_last_modified_date(self._file_path)
        )

    @lazyproperty
    def metadata_file_path(self) -> str | None:
        """The best available file-path for this document or `None` if unavailable."""
        return self._file_path

    @property
    def page_number(self) -> int:
        """The current page (slide) number."""
        return self._page_counter

    @lazyproperty
    def picture_partitioner(self) -> AbstractPicturePartitioner:
        """The sub-partitioner to use for PPTX Picture shapes."""
        # -- Note this value has partitioning-run scope. An instance of this options class is
        # -- instantiated once per partitioning run (each document can have different options).
        # -- Because this is a lazyproperty, it is computed only on the first reference. All
        # -- subsequent references during the same partitioning run will get the same value. This
        # -- ensures Picture shapes are processed consistently within a single document. The
        # -- intended use of `register_picture_partitioner()` is that it be called before processing
        # -- any documents, however there's no reason not to make the mechanism robust against
        # -- unintended use.
        return (
            _NullPicturePartitioner
            if self._PicturePartitionerCls is None
            else self._PicturePartitionerCls
        )

    @lazyproperty
    def pptx_file(self) -> str | IO[bytes]:
        """The PowerPoint document file to be partitioned.

        This is either a str path or a file-like object. `python-pptx` accepts either for opening a
        presentation file.
        """
        if self._file_path:
            return self._file_path

        # -- In Python <3.11 SpooledTemporaryFile does not implement ".seekable" which triggers an
        # -- exception when Zipfile tries to open it. The pptx format is a zip archive so we need
        # -- to work around that bug here.
        if isinstance(self._file, SpooledTemporaryFile):
            self._file.seek(0)
            return io.BytesIO(self._file.read())

        if self._file:
            return self._file

        raise ValueError(
            "No PPTX document specified, either `filename` or `file` argument must be provided"
        )

    @lazyproperty
    def strategy(self) -> str:
        """The requested partitioning strategy.

        This indicates whether the partitioner should undertake expensive operations like inference
        and OCR to produce a more thorough and/or accurate partitioning of the document.

        Can take several values but for PPTX purposes there is only "hi_res" and not "hi_res".
        Depending on the picture-partitioner used, images may only be OCR'ed and added to the
        element-stream when this partitioning strategy is "hi_res".
        """
        return self._strategy

    def table_metadata(self, text_as_html: str | None):
        """ElementMetadata instance suitable for use with Table element."""
        element_metadata = ElementMetadata(
            filename=self.metadata_file_path,
            last_modified=self.last_modified,
            page_number=self.page_number,
            text_as_html=text_as_html,
        )
        element_metadata.detection_origin = DETECTION_ORIGIN
        return element_metadata

    def text_metadata(self, category_depth: int = 0) -> ElementMetadata:
        """ElementMetadata instance suitable for use with Text and subtypes."""
        element_metadata = ElementMetadata(
            filename=self.metadata_file_path,
            last_modified=self.last_modified,
            page_number=self.page_number,
            category_depth=category_depth,
        )
        element_metadata.detection_origin = DETECTION_ORIGIN
        return element_metadata


# ================================================================================================
# SUB-PARTITIONERS
# ================================================================================================


class _NullPicturePartitioner:
    """Does not parse the provided Picture element and generates zero elements."""

    @classmethod
    def iter_elements(cls, picture: Picture, opts: PptxPartitionerOptions) -> Iterator[Element]:
        """No-op picture partitioner."""
        return
        yield
