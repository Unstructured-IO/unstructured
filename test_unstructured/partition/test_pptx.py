# pyright: reportPrivateUsage=false

"""Test suite for `unstructured.partition.pptx` module."""

from __future__ import annotations

import hashlib
import io
import pathlib
import tempfile
from typing import Any, Iterator, cast

import pptx
import pytest
from pptx.shapes.picture import Picture
from pptx.util import Inches
from pytest_mock import MockFixture

from test_unstructured.unit_utils import (
    FixtureRequest,
    Mock,
    assert_round_trips_through_JSON,
    example_doc_path,
    function_mock,
    property_mock,
)
from unstructured.chunking.title import chunk_by_title
from unstructured.documents.elements import (
    Element,
    ElementMetadata,
    Image,
    ListItem,
    NarrativeText,
    PageBreak,
    Text,
    Title,
)
from unstructured.partition.pptx import (
    PptxPartitionerOptions,
    _PptxPartitioner,
    partition_pptx,
    register_picture_partitioner,
)

EXPECTED_PPTX_OUTPUT = [
    Title(text="Adding a Bullet Slide"),
    ListItem(text="Find the bullet slide layout"),
    ListItem(text="Use _TextFrame.text for first bullet"),
    ListItem(text="Use _TextFrame.add_paragraph() for subsequent bullets"),
    NarrativeText(text="Here is a lot of text!"),
    NarrativeText(text="Here is some text in a text box!"),
]


# == document file behaviors =====================================================================


def test_partition_pptx_from_filename():
    elements = partition_pptx(example_doc_path("fake-power-point.pptx"))
    assert elements == EXPECTED_PPTX_OUTPUT
    for element in elements:
        assert element.metadata.filename == "fake-power-point.pptx"


def test_partition_pptx_from_filename_with_metadata_filename():
    elements = partition_pptx(example_doc_path("fake-power-point.pptx"), metadata_filename="test")
    assert elements == EXPECTED_PPTX_OUTPUT
    for element in elements:
        assert element.metadata.filename == "test"


def test_partition_pptx_with_spooled_file():
    """The `partition_pptx() function can handle a `SpooledTemporaryFile.

    Including one that does not have its read-pointer set to the start.
    """
    with open(example_doc_path("fake-power-point.pptx"), "rb") as test_file:
        spooled_temp_file = tempfile.SpooledTemporaryFile()
        spooled_temp_file.write(test_file.read())
        elements = partition_pptx(file=spooled_temp_file)
        assert elements == EXPECTED_PPTX_OUTPUT
        for element in elements:
            assert element.metadata.filename is None


def test_partition_pptx_from_file():
    with open(example_doc_path("fake-power-point.pptx"), "rb") as f:
        elements = partition_pptx(file=f)
    assert elements == EXPECTED_PPTX_OUTPUT
    assert all(e.metadata.filename is None for e in elements)


def test_partition_pptx_from_file_with_metadata_filename():
    with open(example_doc_path("fake-power-point.pptx"), "rb") as f:
        elements = partition_pptx(file=f, metadata_filename="test")
    assert elements == EXPECTED_PPTX_OUTPUT
    for element in elements:
        assert element.metadata.filename == "test"


def test_partition_pptx_raises_with_neither():
    with pytest.raises(ValueError):
        partition_pptx()


def test_partition_pptx_recurses_into_group_shapes():
    elements = partition_pptx(example_doc_path("group-shapes-nested.pptx"))
    assert [e.text for e in elements] == ["A", "B", "C", "D", "E", "F", "G", "H", "I", "J"]


def test_it_loads_a_PPTX_with_a_JPEG_misidentified_as_image_jpg(opts_args: dict[str, Any]):
    opts_args["file_path"] = example_doc_path("test-image-jpg-mime.pptx")
    opts = PptxPartitionerOptions(**opts_args)
    prs = _PptxPartitioner(opts)._presentation
    picture = cast(Picture, prs.slides[0].shapes[0])

    try:
        picture.image
    except AttributeError:
        raise AssertionError("JPEG image not recognized, needs `python-pptx>=1.0.1`")


# == page-break behaviors ========================================================================


def test_partition_pptx_adds_page_breaks(tmp_path: pathlib.Path):
    filename = str(tmp_path / "test-page-breaks.pptx")

    presentation = pptx.Presentation()
    blank_slide_layout = presentation.slide_layouts[6]

    slide = presentation.slides.add_slide(blank_slide_layout)
    left = top = width = height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "This is the first slide."

    slide = presentation.slides.add_slide(blank_slide_layout)
    left = top = width = height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "This is the second slide."

    presentation.save(filename)

    elements = partition_pptx(filename=filename)

    assert elements == [
        NarrativeText(text="This is the first slide."),
        PageBreak(text=""),
        NarrativeText(text="This is the second slide."),
    ]
    for element in elements:
        assert element.metadata.filename == "test-page-breaks.pptx"


def test_partition_pptx_page_breaks_toggle_off(tmp_path: pathlib.Path):
    filename = str(tmp_path / "test-page-breaks.pptx")

    presentation = pptx.Presentation()
    blank_slide_layout = presentation.slide_layouts[6]

    slide = presentation.slides.add_slide(blank_slide_layout)
    left = top = width = height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "This is the first slide."

    slide = presentation.slides.add_slide(blank_slide_layout)
    left = top = width = height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "This is the second slide."

    presentation.save(filename)

    elements = partition_pptx(filename=filename, include_page_breaks=False)

    assert elements == [
        NarrativeText(text="This is the first slide."),
        NarrativeText(text="This is the second slide."),
    ]
    for element in elements:
        assert element.metadata.filename == "test-page-breaks.pptx"


def test_partition_pptx_many_pages():
    elements = partition_pptx(example_doc_path("fake-power-point-many-pages.pptx"))

    # The page_number of PageBreak is None
    assert set(filter(None, (elt.metadata.page_number for elt in elements))) == {1, 2}
    for element in elements:
        assert element.metadata.filename == "fake-power-point-many-pages.pptx"


# == miscellaneous behaviors =====================================================================


def test_partition_pptx_orders_elements(tmp_path: pathlib.Path):
    filename = str(tmp_path / "test-ordering.pptx")
    presentation = pptx.Presentation()
    blank_slide_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank_slide_layout)

    left = top = width = height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "This is lower and should come second"

    left = top = width = height = Inches(1)
    left = top = Inches(-10)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "This is off the page and shouldn't appear"

    left = top = width = height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = ""

    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "This is higher and should come first"

    top = width = height = Inches(1)
    left = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "-------------TOP-------------"

    presentation.save(filename)

    elements = partition_pptx(filename=filename)
    assert elements == [
        Text("-------------TOP-------------"),
        NarrativeText("This is higher and should come first"),
        NarrativeText("This is lower and should come second"),
    ]
    for element in elements:
        assert element.metadata.filename == "test-ordering.pptx"


def test_partition_pptx_grabs_tables():
    elements = partition_pptx(example_doc_path("fake-power-point-table.pptx"))

    assert elements[1].text.startswith("Column 1")
    assert elements[1].text.strip().endswith("Aqua")
    assert elements[1].metadata.text_as_html == (
        "<table>"
        "<tr><td>Column 1</td><td>Column 2</td><td>Column 3</td></tr>"
        "<tr><td>Red</td><td>Green</td><td>Blue</td></tr>"
        "<tr><td>Purple</td><td>Orange</td><td>Yellow</td></tr>"
        "<tr><td>Tangerine</td><td>Pink</td><td>Aqua</td></tr>"
        "</table>"
    )
    assert elements[1].metadata.filename == "fake-power-point-table.pptx"


@pytest.mark.parametrize("infer_table_structure", [True, False])
def test_partition_pptx_infer_table_structure(infer_table_structure: bool):
    elements = partition_pptx(
        example_doc_path("fake-power-point-table.pptx"), infer_table_structure=infer_table_structure
    )
    table_element_has_text_as_html_field = (
        hasattr(elements[1].metadata, "text_as_html")
        and elements[1].metadata.text_as_html is not None
    )
    assert table_element_has_text_as_html_field == infer_table_structure


def test_partition_pptx_malformed():
    elements = partition_pptx(example_doc_path("fake-power-point-malformed.pptx"))

    assert elements[0].text == "Problem Date Placeholder"
    assert elements[1].text == "Test Slide"
    for element in elements:
        assert element.metadata.filename == "fake-power-point-malformed.pptx"


# == image sub-partitioning behaviors ============================================================


def test_partition_pptx_generates_no_Image_elements_by_default():
    assert partition_pptx(example_doc_path("picture.pptx")) == []


def test_partition_pptx_uses_registered_picture_partitioner():
    class FakePicturePartitioner:
        @classmethod
        def iter_elements(cls, picture: Picture, opts: PptxPartitionerOptions) -> Iterator[Element]:
            image_hash = hashlib.sha1(picture.image.blob).hexdigest()
            yield Image(f"Image with hash {image_hash}, strategy: {opts.strategy}")

    register_picture_partitioner(FakePicturePartitioner)

    elements = partition_pptx(example_doc_path("picture.pptx"))

    assert len(elements) == 1
    image = elements[0]
    assert type(image) is Image
    assert image.text == "Image with hash b0a1e6cf904691e6fa42bd9e72acc2b05280dc86, strategy: fast"


# == metadata behaviors ==========================================================================


# -- .metadata.last_modified ---------------------------------------------------------------------


def test_partition_pptx_from_file_path_gets_last_modified_from_filesystem(mocker: MockFixture):
    filesystem_last_modified = "2024-05-01T15:37:28"
    mocker.patch(
        "unstructured.partition.pptx.get_last_modified_date", return_value=filesystem_last_modified
    )

    elements = partition_pptx(example_doc_path("simple.pptx"))

    assert all(e.metadata.last_modified == filesystem_last_modified for e in elements)


def test_partition_pptx_from_file_gets_last_modified_None():
    with open(example_doc_path("simple.pptx"), "rb") as f:
        elements = partition_pptx(file=f)

    assert all(e.metadata.last_modified is None for e in elements)


def test_partition_pptx_from_file_path_prefers_metadata_last_modified(mocker: MockFixture):
    filesystem_last_modified = "2024-05-01T15:37:28"
    metadata_last_modified = "2020-07-05T09:24:28"
    mocker.patch(
        "unstructured.partition.pptx.get_last_modified_date", return_value=filesystem_last_modified
    )

    elements = partition_pptx(
        example_doc_path("simple.pptx"), metadata_last_modified=metadata_last_modified
    )

    assert all(e.metadata.last_modified == metadata_last_modified for e in elements)


def test_partition_pptx_from_file_prefers_metadata_last_modified():
    metadata_last_modified = "2020-07-05T09:24:28"
    with open(example_doc_path("simple.pptx"), "rb") as f:
        elements = partition_pptx(file=f, metadata_last_modified=metadata_last_modified)

    assert all(e.metadata.last_modified == metadata_last_modified for e in elements)


# -- .metadata.languages -------------------------------------------------------------------------


def test_partition_pptx_element_metadata_has_languages():
    elements = partition_pptx(example_doc_path("fake-power-point.pptx"))
    assert elements[0].metadata.languages == ["eng"]


def test_partition_pptx_respects_detect_language_per_element():
    elements = partition_pptx(
        example_doc_path("language-docs/eng_spa_mult.pptx"), detect_language_per_element=True
    )

    langs = [element.metadata.languages for element in elements]
    # languages other than English and Spanish are detected by this partitioner,
    # so this test is slightly different from the other partition tests
    langs = {element.metadata.languages[0] for element in elements if element.metadata.languages}
    assert "eng" in langs
    assert "spa" in langs


def test_partition_pptx_raises_TypeError_for_invalid_languages():
    with pytest.raises(TypeError):
        partition_pptx(example_doc_path("fake-power-point.pptx"), languages="eng")


# == downstream behaviors ========================================================================


def test_partition_pptx_with_json():
    elements = partition_pptx(example_doc_path("fake-power-point.pptx"))
    assert_round_trips_through_JSON(elements)


def test_add_chunking_strategy_by_title_on_partition_pptx():
    filename = example_doc_path("science-exploration-1p.pptx")

    elements = partition_pptx(filename=filename)
    chunk_elements = partition_pptx(filename, chunking_strategy="by_title")
    chunks = chunk_by_title(elements)

    assert chunk_elements != elements
    assert chunk_elements == chunks


def test_partition_pptx_title_shape_detection(tmp_path: pathlib.Path):
    """This tests if the title attribute of a shape is correctly categorized as a title"""
    filename = str(tmp_path / "test-title-shape.pptx")

    # create a fake PowerPoint presentation with a slide containing a title shape
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    assert title_shape is not None
    title_shape.text = (
        "This is a title, it's a bit long so we can make sure it's not narrative text"
    )
    title_shape.text_frame.add_paragraph().text = "this is a subtitle"

    prs.save(filename)

    # partition the PowerPoint presentation and get the first element
    elements = partition_pptx(filename)
    title = elements[0]
    subtitle = elements[1]

    # assert that the first line is a title and has the correct text and depth
    assert isinstance(title, Title)
    assert (
        title.text == "This is a title, it's a bit long so we can make sure it's not narrative text"
    )
    assert title.metadata.category_depth == 0

    # assert that the first line is the subtitle and has the correct text and depth
    assert isinstance(subtitle, Title)
    assert subtitle.text == "this is a subtitle"
    assert subtitle.metadata.category_depth == 1


def test_partition_pptx_level_detection(tmp_path: pathlib.Path):
    """This tests if the level attribute of a paragraph is correctly set as the category depth"""
    filename = str(tmp_path / "test-category-depth.pptx")

    prs = pptx.Presentation()
    blank_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(blank_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    assert title_shape is not None
    title_shape.text = (
        "This is a title, it's a bit long so we can make sure it's not narrative text"
    )
    body_shape = shapes.placeholders[1]

    tf = body_shape.text_frame
    tf.text = "this is the root level bullet"

    p = tf.add_paragraph()
    p.text = "this is the level 1 bullet"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "this is the level 2 bullet"
    p.level = 2

    prs.slides[0].shapes

    prs.save(filename)

    # partition the PowerPoint presentation and get the first element
    elements = partition_pptx(filename)

    # NOTE(newelh) - python_pptx does not create full bullet xml, so unstructured will
    #                not detect the paragraphs as bullets. This is fine for now, as
    #                the level attribute is still set correctly, and what we're testing here
    test_cases = [
        (0, Title, "This is a title, it's a bit long so we can make sure it's not narrative text"),
        (0, NarrativeText, "this is the root level bullet"),
        (1, NarrativeText, "this is the level 1 bullet"),
        (2, NarrativeText, "this is the level 2 bullet"),
    ]

    for element, test_case in zip(elements, test_cases):
        assert element.text == test_case[2], f"expected {test_case[2]}, got {element.text}"
        assert isinstance(
            element,
            test_case[1],
        ), f"expected {test_case[1]}, got {type(element).__name__} for {element.text}"
        assert (
            element.metadata.category_depth == test_case[0]
        ), f"expected {test_case[0]}, got {element.metadata.category_depth} for {element.text}"


def test_partition_pptx_hierarchy_sample_document():
    """This tests if the hierarchy of the sample document is correctly detected"""
    elements = partition_pptx(example_doc_path("sample-presentation.pptx"))

    test_cases = [
        (0, None, "b2859226ba1f9243fb3f1b2ace889f43"),
        (1, "b2859226ba1f9243fb3f1b2ace889f43", "d13f8827e94541c8b818b0df8f942526"),
        (None, None, "cbb95b030de22979af6bfa42969c8202"),
        (0, None, "e535f799d1f0e79d6777efa873a16ce1"),
        (0, "e535f799d1f0e79d6777efa873a16ce1", "f02bbfb417ad60daa2ba35080e96262f"),
        (0, "e535f799d1f0e79d6777efa873a16ce1", "414dfce72ea53cd4649176af0d62a4c1"),
        (1, "414dfce72ea53cd4649176af0d62a4c1", "3d45a95c79473a07db4edca5534a7c49"),
        (1, "414dfce72ea53cd4649176af0d62a4c1", "a33333f527851f700ca175acd04b8a2c"),
        (2, "a33333f527851f700ca175acd04b8a2c", "6f1b87689e4da2b0fb865bc5f92d5702"),
        (0, "e535f799d1f0e79d6777efa873a16ce1", "3f58e0be3b8e8b15cba7adc4eae68586"),
        (None, None, "e5de1b503e64da424fb7d8113371e16d"),
        (0, None, "8319096532fe2e55f66c491ea8313150"),
        (0, "8319096532fe2e55f66c491ea8313150", "17a7e78277ab131a627cb4538bab7390"),
        (0, "8319096532fe2e55f66c491ea8313150", "41a9e1d0390f4edd77181142ceae51bc"),
        (1, "41a9e1d0390f4edd77181142ceae51bc", "cbbc78ef38a035fd66f7b030dcf12f66"),
        (1, "41a9e1d0390f4edd77181142ceae51bc", "2a551e3cbe67561debe0da262a294f24"),
        (2, "2a551e3cbe67561debe0da262a294f24", "7a121a056eedb11ac8804d6fd17afc0c"),
        (0, "8319096532fe2e55f66c491ea8313150", "a24a3caf9853702cb73daae23020b7b4"),
        (0, "8319096532fe2e55f66c491ea8313150", "18367f334b5c8c4602ea413ab68ac35b"),
        (0, "8319096532fe2e55f66c491ea8313150", "7f647b1f0f20c3db40c36ab57d9a5550"),
        (1, "7f647b1f0f20c3db40c36ab57d9a5550", "591c24b41b53aba873188a0881d10961"),
        (1, "7f647b1f0f20c3db40c36ab57d9a5550", "6ec455f5f19782facf184886876c9a66"),
        (2, "6ec455f5f19782facf184886876c9a66", "5614b00c3f6bff23ebba1360e10f6428"),
        (0, "8319096532fe2e55f66c491ea8313150", "2f57a8d4182e6fd5bd5842b0a2d9841b"),
        (None, None, "4120066d251ba675ade42e8a167ca61f"),
        (None, None, "efb9d74b4f8be6308c9a9006da994e12"),
        (0, None, "fd08cacbaddafee5cbacc02528536ee5"),
    ]

    # Zip the test cases with the elements
    for element, test_case in zip(elements, test_cases):
        expected_depth, expected_parent_id, expected_id = test_case
        assert element.metadata.category_depth == expected_depth
        assert element.metadata.parent_id == expected_parent_id
        assert element.id == expected_id


# ================================================================================================
# MODULE-LEVEL FIXTURES
# ================================================================================================


@pytest.fixture()
def opts_args() -> dict[str, Any]:
    """All default arguments for `_XlsxPartitionerOptions`.

    Individual argument values can be changed to suit each test. Makes construction of opts more
    compact for testing purposes.
    """
    return {
        "file": None,
        "file_path": None,
        "include_page_breaks": True,
        "include_slide_notes": False,
        "infer_table_structure": True,
        "strategy": "fast",
    }


# ================================================================================================
# ISOLATED UNIT TESTS
# ================================================================================================
# These test components used by `partition_pptx()` in isolation such that all edge cases can be
# exercised.
# ================================================================================================


class DescribePptxPartitionerOptions:
    """Unit-test suite for `unstructured.partition.xlsx.PptxPartitionerOptions` objects."""

    @pytest.mark.parametrize("arg_value", [True, False])
    def it_knows_whether_to_emit_PageBreak_elements_as_part_of_the_output_element_stream(
        self, arg_value: bool, opts_args: dict[str, Any]
    ):
        opts_args["include_page_breaks"] = arg_value
        opts = PptxPartitionerOptions(**opts_args)

        assert opts.include_page_breaks is arg_value

    @pytest.mark.parametrize("arg_value", [True, False])
    def it_knows_whether_to_partition_content_found_in_slide_notes(
        self, arg_value: bool, opts_args: dict[str, Any]
    ):
        opts_args["include_slide_notes"] = arg_value
        opts = PptxPartitionerOptions(**opts_args)

        assert opts.include_slide_notes is arg_value

    @pytest.mark.parametrize("arg_value", [True, False])
    def it_knows_whether_to_include_text_as_html_in_Table_metadata(
        self, arg_value: bool, opts_args: dict[str, Any]
    ):
        opts_args["infer_table_structure"] = arg_value
        opts = PptxPartitionerOptions(**opts_args)

        assert opts.infer_table_structure is arg_value

    # -- .increment_page_number() ----------------

    def it_generates_a_PageBreak_element_when_the_page_number_is_incremented(
        self, opts_args: dict[str, Any]
    ):
        opts = PptxPartitionerOptions(**opts_args)
        # -- move to the first slide --
        list(opts.increment_page_number())

        page_break_iter = opts.increment_page_number()

        assert isinstance(next(page_break_iter, None), PageBreak)
        assert opts.page_number == 2
        with pytest.raises(StopIteration):
            next(page_break_iter)

    def but_it_does_not_generate_a_PageBreak_element_for_the_first_slide(
        self, opts_args: dict[str, Any]
    ):
        opts = PptxPartitionerOptions(**opts_args)

        page_break_iter = opts.increment_page_number()

        with pytest.raises(StopIteration):
            next(page_break_iter)
        assert opts.page_number == 1

    def and_it_does_not_generate_a_PageBreak_element_when_include_page_breaks_option_is_off(
        self, opts_args: dict[str, Any]
    ):
        opts_args["include_page_breaks"] = False
        opts = PptxPartitionerOptions(**opts_args)
        # -- move to the first slide --
        list(opts.increment_page_number())

        page_break_iter = opts.increment_page_number()

        with pytest.raises(StopIteration):
            next(page_break_iter)
        assert opts.page_number == 2

    # -- .last_modified --------------------------

    def it_gets_last_modified_from_the_filesystem_when_a_path_is_provided(
        self, opts_args: dict[str, Any], get_last_modified_date_: Mock
    ):
        opts_args["file_path"] = "a/b/spreadsheet.pptx"
        get_last_modified_date_.return_value = "2024-04-02T20:32:35"
        opts = PptxPartitionerOptions(**opts_args)

        last_modified = opts.last_modified

        get_last_modified_date_.assert_called_once_with("a/b/spreadsheet.pptx")
        assert last_modified == "2024-04-02T20:32:35"

    def and_it_falls_back_to_None_for_the_last_modified_date_when_no_path_is_provided(
        self, opts_args: dict[str, Any]
    ):
        file = io.BytesIO(b"abcdefg")
        opts_args["file"] = file
        opts = PptxPartitionerOptions(**opts_args)

        last_modified = opts.last_modified

        assert last_modified is None

    # -- .metadata_file_path ---------------------

    @pytest.mark.parametrize("file_path", ["u/v/w.pptx", None])
    def it_uses_the_filename_argument_when_provided(
        self, file_path: str | None, opts_args: dict[str, Any]
    ):
        opts_args["file_path"] = file_path
        opts = PptxPartitionerOptions(**opts_args)

        assert opts.metadata_file_path == file_path

    # -- .page_number ----------------------------

    def it_keeps_track_of_the_page_number(self, opts_args: dict[str, Any]):
        """In PPTX, page-number is the slide number."""
        opts = PptxPartitionerOptions(**opts_args)

        assert opts.page_number == 0
        list(opts.increment_page_number())
        assert opts.page_number == 1
        list(opts.increment_page_number())
        assert opts.page_number == 2

    def it_assigns_the_correct_page_number_when_starting_page_number_is_given(
        self, opts_args: dict[str, Any]
    ):
        opts = PptxPartitionerOptions(**opts_args, starting_page_number=3)
        # -- move to the "first" slide --
        list(opts.increment_page_number())

        table_metadata = opts.table_metadata(text_as_html="<table><tr/></table>")
        text_metadata = opts.text_metadata()

        assert isinstance(table_metadata, ElementMetadata)
        assert isinstance(text_metadata, ElementMetadata)
        assert text_metadata.page_number == 3
        assert table_metadata.page_number == 3

    # -- .pptx_file ------------------------------

    def it_uses_the_path_to_open_the_presentation_when_file_path_is_provided(
        self, opts_args: dict[str, Any]
    ):
        opts_args["file_path"] = "l/m/n.pptx"
        opts = PptxPartitionerOptions(**opts_args)

        assert opts.pptx_file == "l/m/n.pptx"

    def and_it_uses_a_BytesIO_file_to_replaces_a_SpooledTemporaryFile_provided(
        self, opts_args: dict[str, Any]
    ):
        spooled_temp_file = tempfile.SpooledTemporaryFile()
        spooled_temp_file.write(b"abcdefg")
        opts_args["file"] = spooled_temp_file
        opts = PptxPartitionerOptions(**opts_args)

        pptx_file = opts.pptx_file

        assert pptx_file is not spooled_temp_file
        assert isinstance(pptx_file, io.BytesIO)
        assert pptx_file.getvalue() == b"abcdefg"

    def and_it_uses_the_provided_file_directly_when_not_a_SpooledTemporaryFile(
        self, opts_args: dict[str, Any]
    ):
        file = io.BytesIO(b"abcdefg")
        opts_args["file"] = file
        opts = PptxPartitionerOptions(**opts_args)

        pptx_file = opts.pptx_file

        assert pptx_file is file
        assert isinstance(pptx_file, io.BytesIO)
        assert pptx_file.getvalue() == b"abcdefg"

    def but_it_raises_ValueError_when_neither_a_file_path_or_file_is_provided(
        self, opts_args: dict[str, Any]
    ):
        opts = PptxPartitionerOptions(**opts_args)

        with pytest.raises(ValueError, match="No PPTX document specified, either `filename` or "):
            opts.pptx_file

    # -- .strategy -------------------------------

    @pytest.mark.parametrize("arg_value", ["fast", "hi_res"])
    def it_knows_which_partitioning_strategy_to_use(
        self, arg_value: str, opts_args: dict[str, Any]
    ):
        opts_args["strategy"] = arg_value
        opts = PptxPartitionerOptions(**opts_args)

        assert opts.strategy == arg_value

    # -- .table_metadata -------------------------

    def it_can_create_table_metadata(
        self, last_modified_prop_: Mock, metadata_file_path_prop_: Mock, opts_args: dict[str, Any]
    ):
        metadata_file_path_prop_.return_value = "d/e/f.pptx"
        last_modified_prop_.return_value = "2024-04-02T19:51:55"
        opts = PptxPartitionerOptions(**opts_args)
        # -- move to the first slide --
        list(opts.increment_page_number())

        metadata = opts.table_metadata(text_as_html="<table><tr/></table>")

        assert isinstance(metadata, ElementMetadata)
        assert metadata.filename == "f.pptx"
        assert metadata.last_modified == "2024-04-02T19:51:55"
        assert metadata.page_number == 1
        assert metadata.text_as_html == "<table><tr/></table>"

    # -- .text_metadata -------------------------

    def it_can_create_text_metadata(
        self, last_modified_prop_: Mock, metadata_file_path_prop_: Mock, opts_args: dict[str, Any]
    ):
        metadata_file_path_prop_.return_value = "d/e/f.pptx"
        last_modified_prop_.return_value = "2024-04-02T19:56:40"
        opts = PptxPartitionerOptions(**opts_args)
        # -- move to the first slide --
        list(opts.increment_page_number())

        metadata = opts.text_metadata(category_depth=2)

        assert isinstance(metadata, ElementMetadata)
        assert metadata.filename == "f.pptx"
        assert metadata.last_modified == "2024-04-02T19:56:40"
        assert metadata.page_number == 1
        assert metadata.category_depth == 2

    # -- fixtures --------------------------------------------------------------------------------

    @pytest.fixture()
    def get_last_modified_date_(self, request: FixtureRequest):
        return function_mock(request, "unstructured.partition.pptx.get_last_modified_date")

    @pytest.fixture()
    def last_modified_prop_(self, request: FixtureRequest):
        return property_mock(request, PptxPartitionerOptions, "last_modified")

    @pytest.fixture()
    def metadata_file_path_prop_(self, request: FixtureRequest):
        return property_mock(request, PptxPartitionerOptions, "metadata_file_path")
